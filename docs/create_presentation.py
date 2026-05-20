#!/usr/bin/env python3
"""
WSiT Smart IT Support System – Phase 2
Generates: WSiT_Presentation.pptx
Run from the docs/ directory:  python create_presentation.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ─────────────────────────────────────────────────────────────
NAVY  = RGBColor(22,  33,  62)
BLUE  = RGBColor(66, 153, 225)
RED   = RGBColor(233, 69,  96)
GREEN = RGBColor(56, 161, 105)
WHITE = RGBColor(255, 255, 255)
DARK  = RGBColor(45,  55,  72)
GRAY  = RGBColor(113, 128, 150)
LIGHT = RGBColor(237, 242, 247)
AMBER = RGBColor(237, 137,  54)
NAVY2 = RGBColor(30,  45,  85)

# ── Slide dimensions 16:9 ────────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)
M  = Inches(0.4)
CY = Inches(1.28)
CH = SH - CY - Inches(0.25)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def ns():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    r = sl.shapes.add_shape(1, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line: r.line.color.rgb = line
    else: r.line.fill.background()
    return r

def tx(sl, text, l, t, w, h, sz=16, bold=False, col=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    b = sl.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = col
    return b

def hdr(sl, title, sub=None):
    rect(sl, 0, 0, SW, Inches(1.1), NAVY)
    rect(sl, 0, Inches(1.1), SW, Inches(0.06), BLUE)
    tx(sl, title, M, Inches(0.1), SW - M * 2, Inches(0.65),
       sz=24, bold=True, col=WHITE)
    if sub:
        tx(sl, sub, M, Inches(0.72), SW - M * 2, Inches(0.35), sz=12, col=BLUE)

def blist(sl, items, l, t, w, h, sz=16):
    b = sl.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        text = it[0] if isinstance(it, (list, tuple)) else it
        lv   = it[1] if isinstance(it, (list, tuple)) and len(it) > 1 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5 if lv == 0 else 2)
        run = p.add_run()
        run.text = ("• " if lv == 0 else "     – ") + text
        run.font.size = Pt(sz if lv == 0 else sz - 2)
        run.font.color.rgb = DARK if lv == 0 else GRAY
    return b

def card(sl, title, lines, l, t, w, h, bg=LIGHT, tc=NAVY):
    rect(sl, l, t, w, h, bg)
    tx(sl, title, l + Inches(0.13), t + Inches(0.08),
       w - Inches(0.22), Inches(0.35), sz=13, bold=True, col=tc)
    blist(sl, lines, l + Inches(0.13), t + Inches(0.45),
          w - Inches(0.22), h - Inches(0.53), sz=12)

def num(sl, n):
    tx(sl, str(n), SW - Inches(0.65), SH - Inches(0.38),
       Inches(0.5), Inches(0.3), sz=10, col=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# 01 – COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
rect(sl, 0, 0, SW, SH, NAVY)
rect(sl, 0, Inches(4.45), SW, Inches(0.08), BLUE)

tx(sl, "Smart IT Service Request",
   M, Inches(0.75), SW - M * 2, Inches(1.1),
   sz=40, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "& Support Ticket System",
   M, Inches(1.78), SW - M * 2, Inches(1.0),
   sz=40, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "WSiT – World Systems for Information Technology",
   M, Inches(2.85), SW - M * 2, Inches(0.5),
   sz=18, col=BLUE, align=PP_ALIGN.CENTER)
tx(sl, "Software Architecture Project  ·  Phase 1 & 2  ·  May 2026",
   M, Inches(3.38), SW - M * 2, Inches(0.4),
   sz=14, col=GRAY, align=PP_ALIGN.CENTER, italic=True)

rect(sl, Inches(3.5), Inches(4.7), Inches(6.35), Inches(2.45), NAVY2)
tx(sl, "Project Team",
   Inches(3.7), Inches(4.82), Inches(5.9), Inches(0.38),
   sz=13, bold=True, col=BLUE, align=PP_ALIGN.CENTER)
tx(sl, "Zekeriya  ·  Hussein  ·  Hamza  ·  Elsa  ·  Leen",
   Inches(3.5), Inches(5.28), Inches(6.35), Inches(0.45),
   sz=15, col=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "Architecture Style:  Layered REST API  ·  4+1 Model",
   Inches(3.5), Inches(5.85), Inches(6.35), Inches(0.35),
   sz=12, col=GRAY, align=PP_ALIGN.CENTER, italic=True)
tx(sl, "Deadline: 24 May 2026",
   Inches(3.5), Inches(6.28), Inches(6.35), Inches(0.35),
   sz=12, col=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# 02 – AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = ns(); hdr(sl, "Agenda"); num(sl, 2)
blist(sl, [
    "01  Project Overview & Problem Statement",
    "02  Architecture Style & 4+1 Model",
    "03  Technology Stack",
    "04  Use Case View — Actors & Use Cases",
    "05  Logical View — Component Architecture",
    "06  Process View — Workflows & Interactions",
    "07  Development View — Module Structure",
    "08  Deployment View — Dev & Production Topology",
    "09  Phase 2 Key Improvements",
    "10  Team Member Contributions",
    "11  Conclusion",
], M, CY, SW - M * 2, CH, sz=18)


# ══════════════════════════════════════════════════════════════════════════════
# 03 – PROJECT OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns(); hdr(sl, "Project Overview", "WSiT – Smart IT Service Request & Support Ticket System"); num(sl, 3)

card(sl, "The Problem", [
    "IT service requests handled via email & phone",
    "No centralized tracking or ticket management",
    "SLAs missed — no visibility on request status",
    "Technicians had no structured task assignment system",
    "Customers had no way to track their requests",
], M, CY, Inches(5.9), Inches(2.75))

card(sl, "The Solution", [
    "Structured digital ticketing workflow accessible 24/7",
    "Self-service portal for customers (submit & track)",
    "Admin dashboard for full ticket lifecycle management",
    "Technician dashboard for assigned task visibility",
    "Real-time status tracking for all three roles",
], M, CY + Inches(2.88), Inches(5.9), Inches(2.75))

card(sl, "System Users", [
    ("Customer", 0), ("Submits service requests, tracks own tickets", 1),
    ("Technician", 0), ("Views & resolves tickets assigned to them", 1),
    ("Administrator", 0), ("Manages all tickets, assigns technicians", 1),
], Inches(6.5), CY, Inches(6.6), Inches(2.0))

card(sl, "Service Types Supported", [
    "CCTV — Surveillance systems",
    "Networking — LAN/WAN, Wi-Fi",
    "Maintenance — Hardware repair",
    "Smart Home — IoT device installation",
    "Access Control — Biometric & card systems",
    "Tech Support — Software & OS issues",
], Inches(6.5), CY + Inches(2.1), Inches(6.6), Inches(3.53))


# ══════════════════════════════════════════════════════════════════════════════
# 04 – ARCHITECTURE STYLE
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Architecture Style",
    "Layered REST API  ·  4+1 Architectural Model (Kruchten, 1995)")
num(sl, 4)

views = [
    ("Use Case View  (+1)",          "Actors, use cases, and end-to-end scenarios",                      BLUE),
    ("Logical View",                  "Component structure and layer relationships (UML diagrams)",        GREEN),
    ("Process View",                  "Runtime behaviour, workflows, sequence & activity diagrams",       AMBER),
    ("Development View",              "Module organisation, file structure, and dependencies",            RED),
    ("Physical / Deployment View",    "Hardware nodes, Docker containers, deployment topology",           GRAY),
]
for i, (title, desc, color) in enumerate(views):
    y = CY + i * Inches(1.05)
    rect(sl, M, y, Inches(0.18), Inches(0.95), color)
    tx(sl, title, M + Inches(0.28), y + Inches(0.06),
       Inches(5.5), Inches(0.38), sz=14, bold=True, col=NAVY)
    tx(sl, desc,  M + Inches(0.28), y + Inches(0.47),
       Inches(5.5), Inches(0.38), sz=13, col=DARK)

card(sl, "Four-Layer Architecture", [
    ("Presentation Layer", 0), ("React SPA — UI, routing, state management", 1),
    ("Controller Layer", 0),   ("FastAPI routes — request validation & dispatch", 1),
    ("Business Logic Layer", 0), ("Services — rules, JWT auth, CRUD operations", 1),
    ("Data Access Layer", 0),  ("SQLAlchemy ORM — SQLite queries & models", 1),
], Inches(7.1), CY, Inches(6.0), Inches(5.5))


# ══════════════════════════════════════════════════════════════════════════════
# 05 – TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = ns(); hdr(sl, "Technology Stack"); num(sl, 5)

stack = [
    ("Frontend",        "React 18, Vite 5, React Router 6, Axios",      BLUE),
    ("Backend",         "FastAPI 0.104, Uvicorn (ASGI server)",          GREEN),
    ("Database",        "SQLite (file-based), SQLAlchemy 2.0 ORM",       AMBER),
    ("Validation",      "Pydantic 2.7 — request/response schemas",       GRAY),
    ("Auth  (Ph. 2)",   "JWT — python-jose (HS256 signed tokens)",       RED),
    ("Hashing (Ph. 2)", "bcrypt — passlib 1.7.4 + bcrypt 4.0.1",        RED),
    ("Config  (Ph. 2)", "python-dotenv — .env file loading",             NAVY),
    ("Container",       "Docker, docker-compose, Nginx (reverse proxy)", NAVY),
]
cw = (SW - M * 3) / 2
for i, (label, tech, color) in enumerate(stack):
    ci = i % 2; ri = i // 2
    l = M + ci * (cw + M)
    t = CY + ri * Inches(1.05)
    rect(sl, l, t, Inches(0.12), Inches(0.88), color)
    tx(sl, label, l + Inches(0.22), t + Inches(0.06),
       cw - Inches(0.3), Inches(0.35), sz=13, bold=True, col=NAVY)
    tx(sl, tech,  l + Inches(0.22), t + Inches(0.45),
       cw - Inches(0.3), Inches(0.38), sz=13, col=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# 06 – USE CASE VIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Use Case View", "Actors, Use Cases & System Functionality")
num(sl, 6)

card(sl, "Customer", [
    "UC-01  Register Account",
    "UC-02  Login (JWT issued)",
    "UC-03  Submit Service Request",
    "UC-04  View My Tickets",
    "UC-05  View Ticket Details",
    "UC-09  Search & Filter Tickets",
], M, CY, Inches(4.1), CH)

card(sl, "Technician  (Phase 2)", [
    "UC-02  Login (JWT issued)",
    "UC-05  View Ticket Details",
    "UC-10  View Assigned Tickets",
    "UC-11  Update Ticket Status",
    "         Mark Resolved → Closed",
], Inches(4.7), CY, Inches(4.1), CH, tc=GREEN)

card(sl, "Administrator", [
    "UC-02  Login (JWT issued)",
    "UC-05  View Ticket Details",
    "UC-06  View All Tickets",
    "UC-07  Assign Technician (dropdown)",
    "UC-08  Update Any Ticket Status",
    "UC-09  Search & Filter Tickets",
], Inches(9.4), CY, Inches(3.7), CH, tc=RED)


# ══════════════════════════════════════════════════════════════════════════════
# 07 – LOGICAL VIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Logical View", "Component Architecture — Four-Layer Structure")
num(sl, 7)

layers = [
    ("Presentation Layer  (React SPA)", [
        "LoginPage.jsx — Login + Register + Account Type selector",
        "CustomerDashboard.jsx — Submit form, filter bar, ticket list",
        "AdminDashboard.jsx — Search/filter, technician dropdown, manage tickets",
        "TechnicianDashboard.jsx — Assigned tickets, resolve/close buttons",
        "api/client.js — Axios instance with JWT Bearer interceptor & VITE_API_URL",
    ], BLUE),
    ("Controller Layer  (FastAPI Routes)", [
        "routes/auth.py → POST /auth/login,  POST /customers/register,  GET /technicians",
        "routes/tickets.py → /tickets, /tickets/my, /tickets/assigned, /tickets/{id}, /resolve",
    ], GREEN),
    ("Business Logic Layer  (Services)", [
        "services/auth_service.py → JWT create/decode, bcrypt hash/verify, admin seed",
        "services/ticket_service.py → CRUD operations, _apply_filters() search/filter",
    ], AMBER),
    ("Data Access Layer  (SQLAlchemy + Config)", [
        "models/models.py → User (customer/technician/admin), Ticket ORM models + Enums",
        "database/database.py → SQLAlchemy engine, SessionLocal, get_db dependency",
        "core/config.py → Loads SECRET_KEY, ALGORITHM, DATABASE_URL from .env",
    ], RED),
]
for i, (name, items, color) in enumerate(layers):
    t = CY + i * Inches(1.42)
    rect(sl, M, t, SW - M * 2, Inches(1.35), LIGHT)
    rect(sl, M, t, Inches(0.2), Inches(1.35), color)
    tx(sl, name, M + Inches(0.3), t + Inches(0.07),
       Inches(4.5), Inches(0.35), sz=13, bold=True, col=color)
    blist(sl, items, M + Inches(0.3), t + Inches(0.44),
          SW - M * 2 - Inches(0.4), Inches(0.85), sz=12)


# ══════════════════════════════════════════════════════════════════════════════
# 08 – PROCESS VIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Process View", "Key Runtime Workflows & Interactions")
num(sl, 8)

card(sl, "1. JWT Login Flow", [
    "User submits email + password → POST /auth/login",
    "auth_service.verify_password() checks bcrypt hash from DB",
    "auth_service.create_access_token() signs HS256 JWT (24 h TTL)",
    "Frontend stores JWT in localStorage",
    "Axios interceptor attaches  'Authorization: Bearer <JWT>'  to all requests",
], M, CY, Inches(6.4), Inches(2.78))

card(sl, "2. Ticket Submission Flow  (Customer)", [
    "Customer fills form → POST /tickets  with Authorization header",
    "get_current_user() decodes JWT → retrieves User from DB by ID",
    "ticket_service.create_ticket() → INSERT INTO tickets (status = open)",
    "Response returned; ticket appears in customer's filtered list",
], M, CY + Inches(2.92), Inches(6.4), Inches(2.72))

card(sl, "3. Technician Assignment Flow  (Admin)", [
    "Admin opens dashboard → GET /tickets (all tickets)",
    "Admin fetches technician list → GET /technicians (dropdown)",
    "Selects technician, clicks Assign → PUT /tickets/{id}/assign",
    "assigned_technician = name,  status = in_progress",
    "Technician logs in → GET /tickets/assigned (their tasks only)",
], Inches(6.7), CY, Inches(6.4), Inches(2.78))

card(sl, "4. Resolution Flow  (Technician)", [
    "Technician clicks 'Mark Resolved' → PUT /tickets/{id}/resolve",
    "Backend checks: assigned_technician == current_user.full_name",
    "Status updated: in_progress → resolved",
    "Technician clicks 'Mark Closed' → resolved → closed",
    "Customer sees updated status in their dashboard",
], Inches(6.7), CY + Inches(2.92), Inches(6.4), Inches(2.72))


# ══════════════════════════════════════════════════════════════════════════════
# 09 – DEVELOPMENT VIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Development View", "Module Structure & Source Code Organisation")
num(sl, 9)

card(sl, "Backend (FastAPI)", [
    "main.py — FastAPI app, CORS, lifespan, admin seed",
    "app/core/config.py — .env loader (Phase 2)",
    "app/models/models.py — User, Ticket ORM + Enums",
    "app/schemas/schemas.py — Pydantic contracts",
    "app/routes/auth.py — Login, Register, /technicians",
    "app/routes/tickets.py — All endpoints + JWT auth deps",
    "app/services/auth_service.py — JWT + bcrypt + seed",
    "app/services/ticket_service.py — CRUD + filters",
    "Dockerfile — python:3.11-slim image",
    ".env / .env.example — Secrets & config",
], M, CY, Inches(6.3), CH)

card(sl, "Frontend (React + Vite)", [
    "src/App.jsx — Router + RequireAuth + role redirects",
    "src/api/client.js — Axios + JWT + VITE_API_URL",
    "src/pages/LoginPage.jsx — Login/Register + role selector",
    "src/pages/CustomerDashboard.jsx — Submit + filter",
    "src/pages/AdminDashboard.jsx — Manage + assign dropdown",
    "src/pages/TechnicianDashboard.jsx — Tasks (Phase 2)",
    "src/pages/TicketDetails.jsx — Read-only detail view",
    "src/components/Navbar.jsx + StatusBadge.jsx",
    "Dockerfile — Multi-stage: node build → nginx serve",
    ".env / .env.example — VITE_API_URL",
], Inches(6.7), CY, Inches(6.3), CH)


# ══════════════════════════════════════════════════════════════════════════════
# 10 – DEPLOYMENT VIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Deployment View", "Development (Local) & Production (Docker + Nginx)")
num(sl, 10)

card(sl, "Development Deployment  (Local)", [
    "All services run on a single developer machine",
    "Browser → Vite Dev Server :5173 (serves React SPA)",
    "Browser → Uvicorn :8000 (FastAPI REST API)",
    "SQLite file: wsit_support.db on local filesystem",
    ".env file loaded by python-dotenv at startup",
    "CORS allows :5173 and :3000 origins",
], M, CY, Inches(6.3), CH)

card(sl, "Production Deployment  (Docker + Nginx)", [
    "Nginx Reverse Proxy on port 80 (public entry point)",
    ("Route  /  → Frontend container (nginx:alpine)", 1),
    ("Route  /api/*  → Backend container (port 8000)", 1),
    "Frontend Container: React build served by nginx:alpine",
    "Backend Container: python:3.11-slim + Uvicorn",
    "Docker Volume: persistent SQLite database file",
    "docker-compose.yml orchestrates all services",
    "Secrets injected via .env at runtime",
], Inches(6.7), CY, Inches(6.3), CH)


# ══════════════════════════════════════════════════════════════════════════════
# 11 – PHASE 2 IMPROVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Phase 2 — Key Improvements", "Security · New Role · Functionality · Deployment")
num(sl, 11)

improvements = [
    ("JWT Authentication",    "Replaced mock base64 token with signed HS256 JWT via python-jose. Token payload: sub (user ID), role, email, exp.", BLUE),
    ("bcrypt Password Hashing","Replaced SHA-256 with bcrypt (passlib 1.7.4). Each hash is uniquely salted. bcrypt pinned to 4.0.1 for compatibility.", GREEN),
    ("Technician Role",       "New 3rd actor: register as Technician, dedicated dashboard, /tickets/assigned, PUT /tickets/{id}/resolve endpoints.", AMBER),
    ("Search & Filter",       "Query params on GET /tickets and GET /tickets/my: ?search= (ilike), ?status=, ?priority= — combinable.", RED),
    ("Environment Config",    ".env support via python-dotenv (backend: SECRET_KEY, DATABASE_URL) and Vite (frontend: VITE_API_URL).", NAVY),
    ("Docker Support",        "Backend Dockerfile + Frontend multi-stage Dockerfile + docker-compose.yml. Frontend served via nginx:alpine.", GRAY),
]
cw2 = (SW - M * 3) / 2
for i, (title, desc, color) in enumerate(improvements):
    ci = i % 2; ri = i // 2
    l = M + ci * (cw2 + M)
    t = CY + ri * Inches(1.73)
    rect(sl, l, t, Inches(6.3), Inches(1.65), LIGHT)
    rect(sl, l, t, Inches(0.16), Inches(1.65), color)
    tx(sl, title, l + Inches(0.26), t + Inches(0.1),
       Inches(5.8), Inches(0.38), sz=14, bold=True, col=NAVY)
    tx(sl, desc,  l + Inches(0.26), t + Inches(0.52),
       Inches(5.8), Inches(1.0), sz=12, col=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# 12 – SYSTEM FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "System Features Overview", "What Each Role Can Do")
num(sl, 12)

card(sl, "Customer Portal", [
    "Register with Customer account type",
    "Login — redirected to /dashboard automatically",
    "Submit tickets (service type, priority, location)",
    "Filter own tickets by status & priority",
    "View full ticket details & assigned technician",
    "Track ticket through: open → in_progress → resolved → closed",
], M, CY, Inches(4.1), CH, tc=BLUE)

card(sl, "Technician Dashboard  (Phase 2)", [
    "Register with Technician account type",
    "Login — redirected to /technician dashboard",
    "View only tickets assigned to them",
    "Filter by status and priority",
    "Mark tickets Resolved when work is done",
    "Mark tickets Closed after resolution",
    "View full ticket details & customer info",
], Inches(4.7), CY, Inches(4.1), CH, tc=GREEN)

card(sl, "Admin Dashboard", [
    "Seeded admin account: admin@wsit.com",
    "Login — redirected to /admin dashboard",
    "View all tickets with search + filter",
    "Assign technician from registered dropdown",
    "Update any ticket status",
    "Summary cards: Open / In Progress / Resolved / Closed",
    "View full details of any ticket",
], Inches(9.4), CY, Inches(3.7), CH, tc=RED)


# ══════════════════════════════════════════════════════════════════════════════
# 13 – ZEKERIYA
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Zekeriya", "Lead Software Architect & Backend Developer")
num(sl, 13)
rect(sl, M, CY, Inches(0.18), CH, BLUE)
blist(sl, [
    "Defined the overall 4+1 Architectural Model and system design strategy",
    "Designed the Layered REST API Architecture (4-layer separation of concerns)",
    "Set up the FastAPI project structure: routes, services, models, schemas",
    "Implemented SQLAlchemy ORM models: User (3 roles), Ticket + all Enums",
    "Designed the relational database schema (SQLite with SQLAlchemy 2.0)",
    "Implemented all REST API endpoints (9 Phase 1 + 3 Phase 2 endpoints)",
    "Phase 2: Migrated authentication to JWT (python-jose HS256, 24 h TTL)",
    "Phase 2: Replaced SHA-256 with bcrypt password hashing (passlib)",
    "Phase 2: Added Technician role + assigned-ticket endpoints",
    "Phase 2: Implemented environment configuration (app/core/config.py)",
    "Phase 2: Created Backend Dockerfile (python:3.11-slim)",
    "Wrote the technical architecture sections of SAD V1 and V2",
], M + Inches(0.3), CY + Inches(0.05), SW - M * 2 - Inches(0.3), CH, sz=16)


# ══════════════════════════════════════════════════════════════════════════════
# 14 – HUSSEIN
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Hussein", "Frontend Developer")
num(sl, 14)
rect(sl, M, CY, Inches(0.18), CH, GREEN)
blist(sl, [
    "Designed and implemented the full React SPA architecture with Vite 5",
    "Built all 5 application pages: Login, Customer, Admin, Technician, TicketDetails",
    "Implemented role-based route protection using RequireAuth component in App.jsx",
    "LoginPage: Login/Register tabs with account type selector (Customer / Technician)",
    "CustomerDashboard: Ticket submission form, status/priority filter bar, ticket table",
    "AdminDashboard: Search input, filter dropdowns, technician assignment dropdown",
    "Phase 2: Built TechnicianDashboard.jsx — full assigned-ticket management page",
    "Phase 2: Added search/filter UI across all dashboards",
    "Phase 2: Added loading states on all data-fetching components",
    "Configured Axios HTTP client (api/client.js) with JWT Bearer token interceptor",
    "Phase 2: Migrated API base URL to VITE_API_URL environment variable",
    "Created frontend Dockerfile (multi-stage: node:18-alpine → nginx:alpine)",
], M + Inches(0.3), CY + Inches(0.05), SW - M * 2 - Inches(0.3), CH, sz=16)


# ══════════════════════════════════════════════════════════════════════════════
# 15 – HAMZA
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Hamza", "Systems Analyst")
num(sl, 15)
rect(sl, M, CY, Inches(0.18), CH, AMBER)
blist(sl, [
    "Led system requirements gathering, analysis, and validation",
    "Identified and documented all system actors: Customer, Technician, Administrator",
    "Defined 11 use cases (UC-01 to UC-11) with descriptions and actor mappings",
    "Created the Use Case Diagram (PlantUML) for Phase 1 (UC-01 to UC-08)",
    "Analysed system workflows, business rules, and ticket lifecycle transitions",
    "Created the Activity Diagram for Admin Ticket Management workflow",
    "Phase 2: Extended use cases for Technician role (UC-10, UC-11)",
    "Phase 2: Updated Use Case Diagram to include Technician actor",
    "Phase 2: Analysed and documented the technician assignment-to-resolution flow",
    "Contributed to SAD V1 Sections 2 (System Selection) and 3.1 (Use Case View)",
    "Contributed to SAD V2 Sections 3.1 updates and 3.6 Scenarios (+1 View)",
    "Validated all documented use cases against the implemented system",
], M + Inches(0.3), CY + Inches(0.05), SW - M * 2 - Inches(0.3), CH, sz=16)


# ══════════════════════════════════════════════════════════════════════════════
# 16 – ELSA
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Elsa", "API & Integration Specialist")
num(sl, 16)
rect(sl, M, CY, Inches(0.18), CH, RED)
blist(sl, [
    "Designed the complete REST API endpoint specification (methods, schemas, status codes)",
    "Defined Pydantic request/response schemas for all endpoints",
    "Implemented and validated frontend-backend integration across all pages",
    "Tested all API integrations using FastAPI Swagger UI (localhost:8000/docs)",
    "Phase 2: Designed search/filter query parameter API (?status, ?priority, ?search)",
    "Phase 2: Implemented GET /technicians endpoint (admin-accessible technician list)",
    "Phase 2: Implemented GET /tickets/assigned (technician's personal task list)",
    "Phase 2: Implemented PUT /tickets/{id}/resolve (technician status update)",
    "Phase 2: Set up backend .env configuration (SECRET_KEY, DATABASE_URL)",
    "Phase 2: Set up frontend .env configuration (VITE_API_URL)",
    "Wrote the API Endpoint Summary table in README.md",
    "Contributed to SAD Sections 3.2 (Logical View) and 3.3 (Process View)",
], M + Inches(0.3), CY + Inches(0.05), SW - M * 2 - Inches(0.3), CH, sz=16)


# ══════════════════════════════════════════════════════════════════════════════
# 17 – LEEN
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Leen", "QA Engineer & Technical Writer")
num(sl, 17)
rect(sl, M, CY, Inches(0.18), CH, GRAY)
blist(sl, [
    "Wrote and maintained the Software Architecture Document (SAD V1 and V2)",
    "Structured the 4+1 architectural document per Kruchten (1995) methodology",
    "Phase 1 testing: registration, login, ticket submission, admin assign/status flows",
    "Phase 2 testing: JWT authentication and bcrypt password hashing",
    "Phase 2 testing: Search & filter functionality on all ticket endpoints",
    "Phase 2 testing: Full Technician flow (register → login → tasks → resolve → close)",
    "Wrote the README.md with setup instructions for both development and Docker",
    "Documented all API endpoints, environment variables, and filter query parameters",
    "Contributed to SAD Sections 3.4 (Development View) and 3.5 (Deployment View)",
    "Phase 2: Authored the Production Deployment section with Nginx configuration",
    "Validated Docker Compose deployment (docker-compose up --build)",
    "Authored the Conclusion sections for SAD V1 and V2",
], M + Inches(0.3), CY + Inches(0.05), SW - M * 2 - Inches(0.3), CH, sz=16)


# ══════════════════════════════════════════════════════════════════════════════
# 18 – CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sl = ns()
hdr(sl, "Conclusion & Summary")
num(sl, 18)

card(sl, "Phase 1 Deliverables  ✓  (10 Marks)", [
    "SAD Version 1: Use Case View, Logical View, Process View",
    "FastAPI backend: JWT auth, all 9 endpoints, role-based access",
    "React frontend: Login, Customer Dashboard, Admin Dashboard",
    "SQLite database with SQLAlchemy ORM models",
], M, CY, Inches(6.3), Inches(2.82))

card(sl, "Phase 2 Deliverables  ✓  (15 Marks)", [
    "SAD Version 2: + Development View, Deployment View, Phase 2 scenarios",
    "JWT authentication (HS256) + bcrypt password hashing",
    "Technician role: dashboard, /tickets/assigned, resolve endpoints",
    "Search & filter API + UI across all dashboards",
    ".env configuration (backend + frontend)",
    "Docker support: Dockerfile (×2) + docker-compose.yml",
], M, CY + Inches(2.95), Inches(6.3), Inches(2.72))

card(sl, "Architecture Quality", [
    "Layered REST API — clean 4-layer separation of concerns",
    "4+1 Model — all 5 views fully documented (SAD V1 + V2)",
    "Extensible design: Technician role added without redesign",
    "Production-ready: JWT, bcrypt, Docker, Nginx, .env",
    "University-appropriate scope and complexity",
], Inches(6.7), CY, Inches(6.3), Inches(2.82))

rect(sl, Inches(6.7), CY + Inches(2.95), Inches(6.3), Inches(2.72), NAVY)
tx(sl, "Thank You",
   Inches(7.0), CY + Inches(3.35), Inches(5.7), Inches(0.9),
   sz=32, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "WSiT – Smart IT Service Request & Support Ticket System\nSoftware Architecture Project — Phase 1 & 2",
   Inches(7.0), CY + Inches(4.3), Inches(5.7), Inches(1.0),
   sz=13, col=BLUE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).parent / "WSiT_Presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
